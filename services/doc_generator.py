"""Document generator service — generate Word, Excel, PowerPoint from chat content.

Supports:
- Word (.docx) from markdown text
- Excel (.xlsx) from structured data
- PowerPoint (.pptx) from slide content
- PDF from markdown (via HTML conversion)
"""

from __future__ import annotations

import io
import os
import re
import tempfile
import time
from pathlib import Path
from typing import Any

# Output directory
OUTPUT_DIR = Path(os.environ.get("DOC_OUTPUT_DIR", "/app/data/generated_docs"))
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)


def _safe_filename(name: str) -> str:
    """Make a safe filename from a string."""
    name = re.sub(r'[^\w\s\-]', '', name)
    name = re.sub(r'\s+', '_', name.strip())
    return name[:50] or "document"


def generate_word(title: str, content: str) -> dict[str, Any]:
    """Generate a Word document from markdown-like content."""
    from docx import Document
    from docx.shared import Inches, Pt
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc = Document()

    # Style setup
    style = doc.styles['Normal']
    font = style.font
    font.name = 'Calibri'
    font.size = Pt(11)

    # Title
    title_para = doc.add_heading(title, level=0)
    title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER

    # Parse content by lines
    lines = content.split('\n')
    i = 0
    while i < len(lines):
        line = lines[i].strip()

        if not line:
            i += 1
            continue

        # Headings
        if line.startswith('# '):
            doc.add_heading(line[2:], level=1)
        elif line.startswith('## '):
            doc.add_heading(line[3:], level=2)
        elif line.startswith('### '):
            doc.add_heading(line[4:], level=3)
        elif line.startswith('#### '):
            doc.add_heading(line[5:], level=4)
        # Bullet points
        elif line.startswith('- ') or line.startswith('* '):
            doc.add_paragraph(line[2:], style='List Bullet')
        elif re.match(r'^\d+\.\s', line):
            text = re.sub(r'^\d+\.\s', '', line)
            doc.add_paragraph(text, style='List Number')
        # Code blocks
        elif line.startswith('```'):
            code_lines = []
            i += 1
            while i < len(lines) and not lines[i].strip().startswith('```'):
                code_lines.append(lines[i])
                i += 1
            if code_lines:
                code_text = '\n'.join(code_lines)
                p = doc.add_paragraph()
                p.style = doc.styles['Normal']
                run = p.add_run(code_text)
                run.font.name = 'Consolas'
                run.font.size = Pt(9)
                # Light gray background via shading
        # Table (simple markdown table)
        elif '|' in line and line.startswith('|'):
            table_lines = []
            while i < len(lines) and '|' in lines[i]:
                table_lines.append(lines[i])
                i += 1
            i -= 1  # Will be incremented at end of loop
            _add_docx_table(doc, table_lines)
        # Bold line
        elif line.startswith('**') and line.endswith('**'):
            p = doc.add_paragraph()
            run = p.add_run(line.strip('*'))
            run.bold = True
        # Horizontal rule
        elif line.startswith('---') or line.startswith('***'):
            doc.add_paragraph('─' * 40)
        # Regular paragraph
        else:
            doc.add_paragraph(line)

        i += 1

    # Save
    filename = f"{_safe_filename(title)}_{int(time.time())}.docx"
    filepath = OUTPUT_DIR / filename
    doc.save(str(filepath))

    return {
        "ok": True,
        "filename": filename,
        "filepath": str(filepath),
        "format": "docx",
        "size_bytes": filepath.stat().st_size,
    }


def generate_excel(title: str, content: str, tables: list[list[list[str]]] | None = None) -> dict[str, Any]:
    """Generate an Excel file from content. If tables are provided, use them.
    Otherwise, parse markdown tables from content."""
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

    wb = Workbook()

    # Sheet 1: Summary
    ws = wb.active
    ws.title = "Summary"

    # Header style
    header_font = Font(name='Calibri', bold=True, color='FFFFFF', size=11)
    header_fill = PatternFill(start_color='292524', end_color='292524', fill_type='solid')
    header_alignment = Alignment(horizontal='center', vertical='center')
    thin_border = Border(
        left=Side(style='thin'), right=Side(style='thin'),
        top=Side(style='thin'), bottom=Side(style='thin'),
    )

    # Title row
    ws.merge_cells('A1:F1')
    title_cell = ws['A1']
    title_cell.value = title
    title_cell.font = Font(name='Calibri', bold=True, size=14)
    title_cell.alignment = Alignment(horizontal='center')

    # Parse tables from content
    parsed_tables = tables or _parse_markdown_tables(content)

    if parsed_tables:
        for t_idx, table_data in enumerate(parsed_tables):
            start_row = 3 + t_idx * (len(table_data) + 2)

            for r_idx, row in enumerate(table_data):
                for c_idx, cell in enumerate(row):
                    cell_obj = ws.cell(row=start_row + r_idx, column=c_idx + 1, value=cell.strip())
                    if r_idx == 0:  # Header row
                        cell_obj.font = header_font
                        cell_obj.fill = header_fill
                    cell_obj.border = thin_border
                    cell_obj.alignment = Alignment(wrap_text=True)
    else:
        # No tables found, put content as text in column A
        ws.cell(row=3, column=1, value="Content")
        ws.cell(row=3, column=1).font = Font(bold=True)
        for line_idx, line in enumerate(content.split('\n')[:100], start=4):
            ws.cell(row=line_idx, column=1, value=line)

    # Auto-width columns
    for col_idx in range(1, ws.max_column + 1):
        max_length = 0
        for row_idx in range(1, ws.max_row + 1):
            cell = ws.cell(row=row_idx, column=col_idx)
            if cell.value and not hasattr(cell, 'merged_cell'):
                max_length = max(max_length, len(str(cell.value)))
        from openpyxl.utils import get_column_letter
        ws.column_dimensions[get_column_letter(col_idx)].width = min(max_length + 2, 40)

    # Save
    filename = f"{_safe_filename(title)}_{int(time.time())}.xlsx"
    filepath = OUTPUT_DIR / filename
    wb.save(str(filepath))

    return {
        "ok": True,
        "filename": filename,
        "filepath": str(filepath),
        "format": "xlsx",
        "size_bytes": filepath.stat().st_size,
    }


def generate_powerpoint(title: str, content: str) -> dict[str, Any]:
    """Generate a PowerPoint presentation from markdown content."""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title slide
    slide_layout = prs.slide_layouts[0]  # Title slide
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(40)
    title_shape.text_frame.paragraphs[0].font.bold = True

    # Parse content into slides
    lines = content.split('\n')
    current_slide = None
    bullet_items = []

    def _flush_bullets():
        nonlocal bullet_items, current_slide
        if current_slide and bullet_items:
            body = current_slide.placeholders[1]
            tf = body.text_frame
            for item in bullet_items:
                p = tf.add_paragraph()
                p.text = item
                p.font.size = Pt(18)
                p.level = 0
            bullet_items = []

    for line in lines:
        stripped = line.strip()

        if not stripped:
            continue

        # New slide on ## headings
        if stripped.startswith('## '):
            _flush_bullets()
            slide_layout = prs.slide_layouts[1]  # Title + Content
            current_slide = prs.slides.add_slide(slide_layout)
            current_slide.shapes.title.text = stripped[3:]
            bullet_items = []

        # Sub-headings as slide title
        elif stripped.startswith('### '):
            _flush_bullets()
            slide_layout = prs.slide_layouts[1]
            current_slide = prs.slides.add_slide(slide_layout)
            current_slide.shapes.title.text = stripped[4:]
            bullet_items = []

        # Bullet items
        elif stripped.startswith('- ') or stripped.startswith('* '):
            bullet_items.append(stripped[2:])

        # Numbered items
        elif re.match(r'^\d+\.\s', stripped):
            bullet_items.append(re.sub(r'^\d+\.\s', '', stripped))

        # Skip code blocks
        elif stripped.startswith('```'):
            continue

        # Regular text → add as bullet
        elif not stripped.startswith('#') and not stripped.startswith('---'):
            if current_slide is None:
                # First content → create slide
                slide_layout = prs.slide_layouts[1]
                current_slide = prs.slides.add_slide(slide_layout)
                current_slide.shapes.title.text = title
                bullet_items = []
            bullet_items.append(stripped)

    _flush_bullets()

    # If no slides were created beyond title, add a content slide
    if len(prs.slides) < 2:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "Content"
        body = slide.placeholders[1]
        tf = body.text_frame
        for line in content.split('\n')[:20]:
            if line.strip():
                p = tf.add_paragraph()
                p.text = line.strip()
                p.font.size = Pt(16)

    # Save
    filename = f"{_safe_filename(title)}_{int(time.time())}.pptx"
    filepath = OUTPUT_DIR / filename
    prs.save(str(filepath))

    return {
        "ok": True,
        "filename": filename,
        "filepath": str(filepath),
        "format": "pptx",
        "size_bytes": filepath.stat().st_size,
    }


def generate_pdf(title: str, content: str) -> dict[str, Any]:
    """Generate a PDF from markdown content using HTML conversion."""
    import subprocess

    # Convert markdown to HTML
    html_content = _markdown_to_html(title, content)

    # Use wkhtmltopdf if available, else use a simple approach
    html_path = OUTPUT_DIR / f"{_safe_filename(title)}_{int(time.time())}.html"
    html_path.write_text(html_content, encoding='utf-8')

    pdf_path = html_path.with_suffix('.pdf')

    try:
        subprocess.run(
            ['wkhtmltopdf', '--quiet', str(html_path), str(pdf_path)],
            timeout=30, check=True, capture_output=True,
        )
    except (FileNotFoundError, subprocess.TimeoutExpired):
        # Fallback: return HTML as the "PDF" (browser can print it)
        return {
            "ok": True,
            "filename": html_path.name,
            "filepath": str(html_path),
            "format": "html",
            "size_bytes": html_path.stat().st_size,
            "note": "HTML format (install wkhtmltopdf for PDF)",
        }
    finally:
        # Clean up HTML
        if html_path.exists():
            html_path.unlink()

    return {
        "ok": True,
        "filename": pdf_path.name,
        "filepath": str(pdf_path),
        "format": "pdf",
        "size_bytes": pdf_path.stat().st_size,
    }


# ── Helpers ──────────────────────────────────────────────────────────────────

def _parse_markdown_tables(content: str) -> list[list[list[str]]]:
    """Parse markdown tables into list of tables → rows → cells."""
    tables = []
    current_table = []

    for line in content.split('\n'):
        stripped = line.strip()
        if '|' in stripped and stripped.startswith('|'):
            # Skip separator rows (|---|---|)
            if re.match(r'^\|[\s\-|]+\|$', stripped):
                continue
            cells = [c.strip() for c in stripped.split('|')[1:-1]]
            current_table.append(cells)
        else:
            if current_table:
                tables.append(current_table)
                current_table = []

    if current_table:
        tables.append(current_table)

    return tables


def _add_docx_table(doc, table_lines: list[str]) -> None:
    """Add a markdown table to a Word document."""
    from docx.shared import Pt
    from docx.enum.table import WD_TABLE_ALIGNMENT

    rows = []
    for line in table_lines:
        if re.match(r'^\|[\s\-|]+\|$', line):
            continue
        cells = [c.strip() for c in line.split('|')[1:-1]]
        rows.append(cells)

    if not rows:
        return

    table = doc.add_table(rows=len(rows), cols=len(rows[0]))
    table.alignment = WD_TABLE_ALIGNMENT.CENTER

    for r_idx, row in enumerate(rows):
        for c_idx, cell in enumerate(row):
            if c_idx < len(table.columns):
                table.cell(r_idx, c_idx).text = cell
                for paragraph in table.cell(r_idx, c_idx).paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(10)
                        if r_idx == 0:
                            run.font.bold = True


def _markdown_to_html(title: str, content: str) -> str:
    """Convert markdown content to a simple HTML page."""
    html_lines = [
        '<!DOCTYPE html>',
        '<html><head><meta charset="utf-8">',
        f'<title>{title}</title>',
        '<style>',
        'body { font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", Roboto, sans-serif; max-width: 800px; margin: 40px auto; padding: 0 20px; line-height: 1.6; color: #1c1917; }',
        'h1 { border-bottom: 2px solid #1c1917; padding-bottom: 8px; }',
        'h2 { color: #44403c; margin-top: 24px; }',
        'h3 { color: #57534e; }',
        'table { border-collapse: collapse; width: 100%; margin: 16px 0; }',
        'th, td { border: 1px solid #d6d3d1; padding: 8px 12px; text-align: left; }',
        'th { background-color: #292524; color: white; }',
        'tr:nth-child(even) { background-color: #f5f5f4; }',
        'code { background: #f5f5f4; padding: 2px 6px; border-radius: 4px; font-family: monospace; }',
        'pre { background: #1c1917; color: #e7e5e4; padding: 16px; border-radius: 8px; overflow-x: auto; }',
        'pre code { background: none; color: inherit; }',
        'blockquote { border-left: 4px solid #78716c; margin: 16px 0; padding: 8px 16px; background: #fafaf9; }',
        'ul, ol { padding-left: 24px; }',
        'li { margin: 4px 0; }',
        '</style></head><body>',
    ]

    for line in content.split('\n'):
        stripped = line.strip()
        if not stripped:
            html_lines.append('<br>')
        elif stripped.startswith('# '):
            html_lines.append(f'<h1>{stripped[2:]}</h1>')
        elif stripped.startswith('## '):
            html_lines.append(f'<h2>{stripped[3:]}</h2>')
        elif stripped.startswith('### '):
            html_lines.append(f'<h3>{stripped[4:]}</h3>')
        elif stripped.startswith('- ') or stripped.startswith('* '):
            html_lines.append(f'<li>{stripped[2:]}</li>')
        elif stripped.startswith('```'):
            html_lines.append('<pre><code>')
        elif stripped.startswith('---'):
            html_lines.append('<hr>')
        elif stripped.startswith('> '):
            html_lines.append(f'<blockquote>{stripped[2:]}</blockquote>')
        elif '|' in stripped and stripped.startswith('|'):
            cells = [c.strip() for c in stripped.split('|')[1:-1]]
            tag = 'th' if not re.match(r'^\|[\s\-|]+\|$', stripped) else 'td'
            if not re.match(r'^\|[\s\-|]+\|$', stripped):
                row = ''.join(f'<{tag}>{c}</{tag}>' for c in cells)
                html_lines.append(f'<tr>{row}</tr>')
        else:
            # Inline formatting
            text = stripped
            text = re.sub(r'\*\*(.+?)\*\*', r'<strong>\1</strong>', text)
            text = re.sub(r'\*(.+?)\*', r'<em>\1</em>', text)
            text = re.sub(r'`(.+?)`', r'<code>\1</code>', text)
            html_lines.append(f'<p>{text}</p>')

    html_lines.extend(['</body></html>'])
    return '\n'.join(html_lines)
